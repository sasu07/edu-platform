"""
Generează prezentarea EtoX Platform pentru părinți.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Culori brand ────────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x1E, 0x40, 0xAF)   # #1E40AF
BLUE_MID    = RGBColor(0x25, 0x63, 0xEB)   # #2563EB
BLUE_LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)   # #EFF6FF
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED
PURPLE_LIGHT= RGBColor(0xF5, 0xF3, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)   # #16A34A
GREEN_LIGHT = RGBColor(0xF0, 0xFD, 0xF4)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK   = RGBColor(0x11, 0x18, 0x27)
GRAY_MID    = RGBColor(0x37, 0x41, 0x51)
GRAY_LIGHT  = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_BORDER = RGBColor(0xE5, 0xE7, 0xEB)

# Slide dimensions: widescreen 16:9
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # complet gol

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, border_rgb=None, border_pt=0, radius=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if border_rgb and border_pt:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_text_multiline(slide, lines, x, y, w, h,
                       font_size=16, bold_first=False, color=GRAY_MID,
                       align=PP_ALIGN.LEFT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def slide_header(slide, title, subtitle=None,
                 bg=BLUE_DARK, title_color=WHITE, sub_color=WHITE):
    """Header bar la top."""
    add_rect(slide, 0, 0, W, Inches(1.35), fill_rgb=bg)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                 font_size=16, bold=False, color=sub_color, align=PP_ALIGN.LEFT)

def card(slide, x, y, w, h, icon, title, body_lines,
         accent=BLUE_MID, bg=WHITE, title_size=16, body_size=13):
    """Card cu iconiță, titlu și bullet-uri."""
    add_rect(slide, x, y, w, h, fill_rgb=bg, border_rgb=GRAY_BORDER, border_pt=0.75)
    # accent bar stânga
    add_rect(slide, x, y + Inches(0.12), Inches(0.06), h - Inches(0.24), fill_rgb=accent)
    # icon
    add_text(slide, icon, x + Inches(0.15), y + Inches(0.08), Inches(0.55), Inches(0.55),
             font_size=22, align=PP_ALIGN.CENTER)
    # title
    add_text(slide, title,
             x + Inches(0.72), y + Inches(0.1), w - Inches(0.85), Inches(0.38),
             font_size=title_size, bold=True, color=GRAY_DARK)
    # bullets
    bullet_y = y + Inches(0.48)
    for line in body_lines:
        add_text(slide, f"• {line}",
                 x + Inches(0.72), bullet_y, w - Inches(0.85), Inches(0.3),
                 font_size=body_size, color=GRAY_MID)
        bullet_y += Inches(0.28)

def flow_step(slide, x, y, w, h, num, title, desc, accent=BLUE_MID):
    """Step numbered pentru flow diagrams."""
    add_rect(slide, x, y, w, h, fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
    # cerc număr
    circle = slide.shapes.add_shape(9, x + Inches(0.12), y + Inches(0.12),
                                     Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num)
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    add_text(slide, title, x + Inches(0.65), y + Inches(0.1), w - Inches(0.78), Inches(0.3),
             font_size=13, bold=True, color=GRAY_DARK)
    add_text(slide, desc, x + Inches(0.65), y + Inches(0.38), w - Inches(0.78), Inches(0.55),
             font_size=11, color=GRAY_MID)

def badge(slide, x, y, text, bg=BLUE_MID, text_color=WHITE, size=12):
    bw = Inches(1.6); bh = Inches(0.32)
    add_rect(slide, x, y, bw, bh, fill_rgb=bg)
    add_text(slide, text, x, y, bw, bh, font_size=size, bold=True,
             color=text_color, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 1 — TITLU
# =============================================================================
slide1 = prs.slides.add_slide(BLANK)

# gradient background
add_rect(slide1, 0, 0, W, H, fill_rgb=BLUE_DARK)
add_rect(slide1, 0, Inches(5.2), W, Inches(2.3), fill_rgb=RGBColor(0x1E, 0x29, 0x5C))

# decorative circles
for (cx, cy, cr, alpha) in [
    (Inches(11.5), Inches(0.5), Inches(2.8), RGBColor(0x25,0x63,0xEB)),
    (Inches(12.5), Inches(2.0), Inches(1.6), RGBColor(0x1D,0x4E,0xD8)),
    (Inches(1.0),  Inches(6.2), Inches(1.4), RGBColor(0x1D,0x4E,0xD8)),
]:
    c = slide1.shapes.add_shape(9, cx - cr/2, cy - cr/2, cr, cr)
    c.fill.solid(); c.fill.fore_color.rgb = alpha
    c.line.fill.background()

# Logo / titlu platformă
add_text(slide1, "EtoX",
         Inches(0.7), Inches(1.4), Inches(6), Inches(1.5),
         font_size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide1, "Platform",
         Inches(0.7), Inches(2.7), Inches(6), Inches(1.0),
         font_size=48, bold=False, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.LEFT)
add_text(slide1, "Platforma digitală de pregătire pentru BAC",
         Inches(0.7), Inches(3.65), Inches(8), Inches(0.6),
         font_size=22, bold=False, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.LEFT)

# Tagline
add_text(slide1,
         "Exerciții · Sesiuni de studiu · Monitorizare · Ajutor live",
         Inches(0.7), Inches(4.35), Inches(10), Inches(0.5),
         font_size=16, italic=True, color=RGBColor(0x7D,0xD3,0xFC), align=PP_ALIGN.LEFT)

# Bottom bar
add_rect(slide1, 0, Inches(6.8), W, Inches(0.7), fill_rgb=RGBColor(0x1D,0x4E,0xD8))
add_text(slide1, "Pregătit de EtoX Academy  •  2026",
         Inches(0.5), Inches(6.85), Inches(8), Inches(0.4),
         font_size=14, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.LEFT)
add_text(slide1, "Prezentare pentru părinți",
         Inches(8.5), Inches(6.85), Inches(4.5), Inches(0.4),
         font_size=14, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.RIGHT)

# =============================================================================
# SLIDE 2 — DE CE ETOX?
# =============================================================================
slide2 = prs.slides.add_slide(BLANK)
add_rect(slide2, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide2, "De ce EtoX?", "O platformă construită pentru succesul elevului la BAC")

# 4 carduri mari
cards = [
    ("🎯", "Exerciții BAC autentice",
     ["Baza de date cu exerciții din examene reale",
      "Filtrate pe subiect, dificultate și profil",
      "Actualizate continuu cu subiecte noi"], BLUE_MID),
    ("👨‍🏫", "Profesor uman, nu AI",
     ["Soluțiile elevilor sunt verificate de profesori",
      "Feedback personalizat pe fiecare soluție",
      "Sesiuni live la cerere pentru blocaje"], GREEN),
    ("📊", "Monitorizare în timp real",
     ["Părintele vede tot progresul copilului",
      "Grafice de activitate zilnică",
      "Notificări la fiecare etapă importantă"], PURPLE),
    ("🏆", "Sistem de motivare",
     ["XP, nivele și badge-uri pentru progres",
      "Streak zilnic — recompensează constanța",
      "Plan săptămânal cu countdown la BAC"], AMBER),
]
for i, (icon, title, bullets, accent) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.35)
    y = Inches(1.55) + row * Inches(2.55)
    card(slide2, x, y, Inches(6.1), Inches(2.35), icon, title, bullets,
         accent=accent, title_size=17, body_size=13)

# =============================================================================
# SLIDE 3 — FLUXUL ELEVULUI
# =============================================================================
slide3 = prs.slides.add_slide(BLANK)
add_rect(slide3, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide3, "Fluxul elevului", "Cum lucrează un elev în platformă — pas cu pas")

steps = [
    ("Generează exerciții",   "Selectează subiectul, dificultatea și profileul dorit. Platforma generează un set personalizat."),
    ("Rezolvă exercițiul",    "Lucrează pe foaie. Poate vedea rezolvarea pentru a verifica metoda."),
    ("Trimite soluția",       "Fotografiază sau scanează foaia și o încarcă în platformă pentru verificare."),
    ("Primește feedback",     "Profesorul corectează și marchează soluția. Dacă e corectă, exercițiul devine rezolvat."),
    ("Câștigă XP și avansezi","Fiecare exercițiu rezolvat aduce XP. Crești în nivel: Debutant → Aspirant → Elite → BAC Ready."),
    ("Planifică sesiuni",     "Setează sesiuni de studiu săptămânale (10 sau 25 exerciții) și urmărești countdown-ul la BAC."),
]

SW = Inches(3.9)
SH = Inches(1.05)
for i, (title, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.55) + row * (SH + Inches(0.15))
    flow_step(slide3, x, y, SW + Inches(2.55), SH, i+1, title, desc, BLUE_MID)

# Arrow between col 1 → col 2 (conceptual, skip for simplicity)

# =============================================================================
# SLIDE 4 — FLUXUL PROFESORULUI
# =============================================================================
slide4 = prs.slides.add_slide(BLANK)
add_rect(slide4, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide4, "Fluxul profesorului", "Cum interacționează profesorul cu elevii în platformă")

teacher_steps = [
    ("Primește notificare",    "Imediat ce un elev încarcă o soluție, profesorul este notificat în platformă."),
    ("Preia submisiile",       "Poate prelua toate submisiile in asteptare dintr-un singur click - Preia toate."),
    ("Corectează soluția",     "Vede soluția elevului, poate adăuga o notă și marchează: Corect sau Incorect."),
    ("Încarcă fișier",         "Poate atașa un PDF sau imagine cu rezolvarea detaliată pentru elev."),
    ("Răspunde la cereri live","Când un elev cere ajutor live (premium), profesorul setează o dată/oră și un link Zoom."),
    ("Planifică pentru elevi", "Profesorul poate adăuga sesiuni de studiu în calendarul oricărui elev."),
]

for i, (title, desc) in enumerate(teacher_steps):
    col = i % 2
    row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.55) + row * (Inches(1.05) + Inches(0.15))
    flow_step(slide4, x, y, Inches(6.3), Inches(1.05), i+1, title, desc, GREEN)

# =============================================================================
# SLIDE 5 — FLUXUL PĂRINTELUI
# =============================================================================
slide5 = prs.slides.add_slide(BLANK)
add_rect(slide5, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide5, "Fluxul părintelui", "Vizibilitate completă asupra activității copilului")

# Stânga — ce vede
add_rect(slide5, Inches(0.35), Inches(1.55), Inches(5.9), Inches(5.55),
         fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
add_rect(slide5, Inches(0.35), Inches(1.55), Inches(5.9), Inches(0.5), fill_rgb=PURPLE)
add_text(slide5, "📊 Ce vede părintele",
         Inches(0.5), Inches(1.57), Inches(5.6), Inches(0.44),
         font_size=15, bold=True, color=WHITE)

parent_sees = [
    ("📅", "Activitate zilnică",      "Grafic cu exerciții rezolvate în ultimele 30 de zile"),
    ("📚", "Progres pe subiecte",      "Câte exerciții rezolvate din S1, S2, S3"),
    ("⚡", "Sesiuni de studiu",        "Ultimele sesiuni: tip, durată, % completare, XP câștigat"),
    ("🏆", "Nivel și XP",             "Nivelul actual al elevului și progresul spre BAC Ready"),
    ("📹", "Sesiuni live programate", "Vede când are elevul sesiune live cu profesorul"),
    ("🔔", "Notificări",              "Primește notificare când e programată o sesiune live"),
]
for i, (icon, title, desc) in enumerate(parent_sees):
    y = Inches(2.2) + i * Inches(0.78)
    add_text(slide5, icon, Inches(0.55), y, Inches(0.4), Inches(0.4), font_size=18)
    add_text(slide5, title, Inches(1.0), y, Inches(2.0), Inches(0.3),
             font_size=13, bold=True, color=GRAY_DARK)
    add_text(slide5, desc, Inches(1.0), y + Inches(0.28), Inches(4.9), Inches(0.35),
             font_size=11, color=GRAY_MID)

# Dreapta — cum se configurează
add_rect(slide5, Inches(6.6), Inches(1.55), Inches(6.38), Inches(5.55),
         fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
add_rect(slide5, Inches(6.6), Inches(1.55), Inches(6.38), Inches(0.5), fill_rgb=PURPLE)
add_text(slide5, "🔗 Cum funcționează legătura",
         Inches(6.75), Inches(1.57), Inches(6.1), Inches(0.44),
         font_size=15, bold=True, color=WHITE)

setup_steps = [
    ("1", 'Elevul intra in contul sau si acceseaza sectiunea "Parintii mei"'),
    ("2", "Introduce adresa de email a parintelui si salveaza"),
    ("3", "Parintele primeste acces automat la dashboard-ul de progres"),
    ("4", "Legatura poate fi eliminata oricand de elev"),
]
for i, (num, desc) in enumerate(setup_steps):
    y = Inches(2.25) + i * Inches(0.9)
    c = slide5.shapes.add_shape(9, Inches(6.85), y, Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = PURPLE
    c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num; r.font.bold = True; r.font.size = Pt(14)
    r.font.color.rgb = WHITE; r.font.name = "Calibri"
    add_text(slide5, desc, Inches(7.4), y + Inches(0.02), Inches(5.3), Inches(0.6),
             font_size=13, color=GRAY_MID)

add_rect(slide5, Inches(6.6), Inches(5.6), Inches(6.38), Inches(1.5),
         fill_rgb=PURPLE_LIGHT, border_rgb=PURPLE, border_pt=0.5)
add_text(slide5, "💡 Confidențialitate",
         Inches(6.8), Inches(5.7), Inches(6.0), Inches(0.35),
         font_size=13, bold=True, color=PURPLE)
add_text(slide5, "Părintele vede DOAR statistici și activitate — nu are acces la exerciții, soluții sau conversații cu profesorul.",
         Inches(6.8), Inches(6.05), Inches(5.9), Inches(0.85),
         font_size=11, color=GRAY_MID)

# =============================================================================
# SLIDE 6 — SISTEM DE GAMIFICARE
# =============================================================================
slide6 = prs.slides.add_slide(BLANK)
add_rect(slide6, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide6, "Motivare prin gamificare", "XP, nivele și badge-uri care fac studiul captivant")

# Nivele
levels = [
    ("⭐", "Debutant",    "0 XP",     GRAY_MID,   "#6B7280"),
    ("🔥", "Aspirant",   "500 XP",   BLUE_MID,   "#2563EB"),
    ("⚡", "Competitor", "1.500 XP", GREEN,       "#16A34A"),
    ("💎", "Elite",      "4.000 XP", PURPLE,      "#7C3AED"),
    ("🏆", "BAC Ready",  "8.000 XP", AMBER,       "#D97706"),
]

add_text(slide6, "Nivelele de progres",
         Inches(0.5), Inches(1.55), Inches(6), Inches(0.38),
         font_size=15, bold=True, color=GRAY_DARK)

LW = Inches(2.35)
for i, (icon, name, xp, color, hex_c) in enumerate(levels):
    x = Inches(0.4) + i * (LW + Inches(0.18))
    y = Inches(2.0)
    lh = Inches(1.3) + i * Inches(0.12)
    add_rect(slide6, x, Inches(2.0) + (Inches(1.9) - lh), LW, lh,
             fill_rgb=color, border_rgb=None)
    add_text(slide6, icon + "  " + name,
             x, Inches(2.0) + (Inches(1.9) - lh) - Inches(0.45), LW, Inches(0.38),
             font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide6, xp,
             x, Inches(2.0) + (Inches(1.9) - lh) - Inches(0.75), LW, Inches(0.3),
             font_size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)

# Arrow progress bar
add_rect(slide6, Inches(0.4), Inches(3.97), Inches(12.15), Inches(0.1), fill_rgb=BLUE_MID)

# Cum se câștigă XP
add_text(slide6, "Cum se câștigă XP?",
         Inches(0.5), Inches(4.25), Inches(6), Inches(0.38),
         font_size=15, bold=True, color=GRAY_DARK)

xp_sources = [
    ("📤", "Trimitere soluție",   "Elev încarcă foto soluție → intră în verificare"),
    ("✅", "Aprobare profesor",   "Profesorul marchează soluția Corectă → XP complet"),
    ("⚡", "Sesiune completată",  "+20 XP (test scurt) sau +50 XP (test BAC ≥80%)"),
    ("🔥", "Streak zilnic",       "Bonus +20% XP la exerciții dacă streak ≥ 3 zile"),
]
for i, (icon, title, desc) in enumerate(xp_sources):
    col = i % 2; row = i // 2
    x = Inches(0.4) + col * Inches(6.45)
    y = Inches(4.75) + row * Inches(1.15)
    add_rect(slide6, x, y, Inches(6.1), Inches(1.0),
             fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
    add_text(slide6, icon, x + Inches(0.12), y + Inches(0.15), Inches(0.5), Inches(0.5), font_size=22)
    add_text(slide6, title, x + Inches(0.7), y + Inches(0.08), Inches(5.1), Inches(0.32),
             font_size=13, bold=True, color=GRAY_DARK)
    add_text(slide6, desc, x + Inches(0.7), y + Inches(0.4), Inches(5.1), Inches(0.45),
             font_size=11, color=GRAY_MID)

# =============================================================================
# SLIDE 7 — SESIUNI DE STUDIU
# =============================================================================
slide7 = prs.slides.add_slide(BLANK)
add_rect(slide7, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide7, "Sesiuni de studiu", "Structura eficientă a muncii — test scurt sau test BAC")

# 2 tipuri de sesiuni
for i, (icon, title, details, accent) in enumerate([
    ("⚡", "Test Scurt", [
        "10 exerciții selectate aleatoriu",
        "Limită de timp: 60 minute",
        "Timer live vizibil pe ecran",
        "Bonus XP: +20 la finalizare",
        "Ideal pentru antrenament zilnic",
    ], BLUE_MID),
    ("🏆", "Test BAC", [
        "~25 exerciții complete (similar examen real)",
        "Limită de timp: 3 ore",
        "Simulare condiții de examen",
        "Bonus XP: +50 dacă ≥80% completare",
        "Ideal pentru pregătire finală",
    ], PURPLE),
]):
    x = Inches(0.4) + i * Inches(6.45)
    add_rect(slide7, x, Inches(1.55), Inches(6.1), Inches(4.2),
             fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
    add_rect(slide7, x, Inches(1.55), Inches(6.1), Inches(0.65), fill_rgb=accent)
    add_text(slide7, icon + "  " + title,
             x + Inches(0.2), Inches(1.6), Inches(5.6), Inches(0.5),
             font_size=22, bold=True, color=WHITE)
    for j, d in enumerate(details):
        add_text(slide7, "✓  " + d,
                 x + Inches(0.25), Inches(2.35) + j * Inches(0.58), Inches(5.5), Inches(0.45),
                 font_size=14, color=GRAY_DARK)

# Rezumat sesiune — bottom
add_rect(slide7, Inches(0.4), Inches(5.9), Inches(12.55), Inches(1.25),
         fill_rgb=BLUE_LIGHT, border_rgb=BLUE_MID, border_pt=0.75)
add_text(slide7, "📈  La finalul fiecărei sesiuni elevul vede:",
         Inches(0.65), Inches(5.98), Inches(5), Inches(0.35),
         font_size=13, bold=True, color=BLUE_DARK)
for i, txt in enumerate(["Exerciții rezolvate / total", "Timp petrecut", "% completare", "XP câștigat + bonus"]):
    add_text(slide7, "• " + txt,
             Inches(0.65) + i * Inches(3.05), Inches(6.35), Inches(2.9), Inches(0.5),
             font_size=12, color=GRAY_MID)

# =============================================================================
# SLIDE 8 — PLAN SĂPTĂMÂNAL & CALENDAR
# =============================================================================
slide8 = prs.slides.add_slide(BLANK)
add_rect(slide8, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide8, "Plan săptămânal & Calendar", "Organizare și disciplină în pregătire")

# Calendar visual (7 coloane)
DAYS = ["Lun", "Mar", "Mie", "Joi", "Vin", "Sam", "Dum"]
CW = Inches(1.68); CH = Inches(2.5)
for i, day in enumerate(DAYS):
    x = Inches(0.35) + i * (CW + Inches(0.06))
    y = Inches(1.6)
    is_today = (i == 1)
    bg = BLUE_LIGHT if is_today else WHITE
    border = BLUE_MID if is_today else GRAY_BORDER
    add_rect(slide8, x, y, CW, CH, fill_rgb=bg, border_rgb=border, border_pt=1.0 if is_today else 0.75)
    add_text(slide8, day, x, y + Inches(0.08), CW, Inches(0.3),
             font_size=11, bold=True, color=BLUE_DARK if is_today else GRAY_MID,
             align=PP_ALIGN.CENTER)
    # Exemple de intrări în calendar
    if i == 1:
        add_rect(slide8, x + Inches(0.08), y + Inches(0.45), CW - Inches(0.16), Inches(0.42),
                 fill_rgb=BLUE_MID)
        add_text(slide8, "⚡ Test Scurt", x + Inches(0.08), y + Inches(0.47),
                 CW - Inches(0.16), Inches(0.35), font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i == 3:
        add_rect(slide8, x + Inches(0.08), y + Inches(0.45), CW - Inches(0.16), Inches(0.42),
                 fill_rgb=PURPLE)
        add_text(slide8, "📹 Live cu prof", x + Inches(0.08), y + Inches(0.47),
                 CW - Inches(0.16), Inches(0.35), font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i == 5:
        add_rect(slide8, x + Inches(0.08), y + Inches(0.45), CW - Inches(0.16), Inches(0.42),
                 fill_rgb=RGBColor(0x1E,0x40,0xAF))
        add_text(slide8, "🏆 Test BAC", x + Inches(0.08), y + Inches(0.47),
                 CW - Inches(0.16), Inches(0.35), font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Features dreapta
features8 = [
    ("📅", "Planificare săptămânală", "Elevul sau profesorul adaugă sesiuni în calendar."),
    ("⏳", "Countdown BAC",           "Numărătoare inversă live până la data BAC 2026."),
    ("👨‍🏫", "Sesiuni de la profesor", "Profesorul poate programa sesiuni direct în calendarul elevului."),
    ("📹", "Ajutor live în calendar", "Sesiunile live apar vizibil diferențiat (mov) în calendar."),
]
for i, (icon, title, desc) in enumerate(features8):
    y = Inches(4.28) + i * Inches(0.77)
    add_text(slide8, icon, Inches(0.5), y, Inches(0.4), Inches(0.38), font_size=18)
    add_text(slide8, title, Inches(1.0), y, Inches(4.0), Inches(0.28),
             font_size=13, bold=True, color=GRAY_DARK)
    add_text(slide8, desc, Inches(1.0), y + Inches(0.28), Inches(11.9), Inches(0.35),
             font_size=11, color=GRAY_MID)

# =============================================================================
# SLIDE 9 — AJUTOR LIVE
# =============================================================================
slide9 = prs.slides.add_slide(BLANK)
add_rect(slide9, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide9, "Ajutor live cu profesorul",
             "Când elevul se blochează, profesorul intervine în timp real", bg=PURPLE)

# Flow stânga
flow9 = [
    ("Elevul se blochează",    "Apasă 'M-am blocat' pe exercițiu și descrie problema", BLUE_MID),
    ("Cer ajutor live",        "Cu abonament Premium: trimite cerere de sesiune live", PURPLE),
    ("Profesorul preia",       "Vede cererea în dashboard și o acceptă", GREEN),
    ("Programează sesiunea",   "Setează data, ora și link-ul Zoom/Meet", AMBER),
    ("Elevul e notificat",     "Primește notificare + apare în calendarul său", BLUE_MID),
    ("Părintele e notificat",  "Primește notificare cu data și ora sesiunii", PURPLE),
]

for i, (title, desc, color) in enumerate(flow9):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.6) + row * Inches(1.55)
    add_rect(slide9, x, y, Inches(6.1), Inches(1.4),
             fill_rgb=WHITE, border_rgb=GRAY_BORDER, border_pt=0.75)
    add_rect(slide9, x, y, Inches(0.08), Inches(1.4), fill_rgb=color)
    c9 = slide9.shapes.add_shape(9, x + Inches(0.2), y + Inches(0.17),
                                   Inches(0.42), Inches(0.42))
    c9.fill.solid(); c9.fill.fore_color.rgb = color
    c9.line.fill.background()
    r9 = c9.text_frame.paragraphs[0].add_run()
    r9.text = str(i+1); r9.font.bold = True; r9.font.size = Pt(13)
    r9.font.color.rgb = WHITE; r9.font.name = "Calibri"
    c9.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide9, title, x + Inches(0.72), y + Inches(0.1),
             Inches(5.2), Inches(0.32), font_size=13, bold=True, color=GRAY_DARK)
    add_text(slide9, desc, x + Inches(0.72), y + Inches(0.42),
             Inches(5.2), Inches(0.75), font_size=11, color=GRAY_MID)

# Note premium
add_rect(slide9, Inches(0.35), Inches(6.65), Inches(12.55), Inches(0.6),
         fill_rgb=PURPLE_LIGHT, border_rgb=PURPLE, border_pt=0.5)
add_text(slide9, "🔒  Ajutorul live este disponibil exclusiv abonaților Premium. Abonamentul poate fi activat de administrator sau școală.",
         Inches(0.6), Inches(6.72), Inches(12.1), Inches(0.45),
         font_size=12, color=PURPLE, italic=True)

# =============================================================================
# SLIDE 10 — ABONAMENTE & ACCES
# =============================================================================
slide10 = prs.slides.add_slide(BLANK)
add_rect(slide10, 0, 0, W, H, fill_rgb=GRAY_LIGHT)
slide_header(slide10, "Tipuri de acces", "Ce include fiecare nivel de abonament")

plans = [
    ("Gratuit", [
        "✓  Generare exerciții (limitat)",
        "✓  Vizualizare rezolvări",
        "✓  Autoevaluare soluții",
        "✓  Gamificare (XP, nivele)",
        "✓  Plan săptămânal de bază",
        "✗  Verificare soluții de profesor",
        "✗  Ajutor live cu profesor",
        "✗  Download variante PDF",
    ], GRAY_MID, WHITE),
    ("Premium", [
        "✓  Generare exerciții nelimitată",
        "✓  Verificare soluții de profesor",
        "✓  Feedback personalizat",
        "✓  Ajutor live cu profesor",
        "✓  Download variante BAC PDF",
        "✓  Plan săptămânal complet",
        "✓  Statistici detaliate",
        "✓  Monitorizare pentru părinți",
    ], BLUE_MID, WHITE),
]

PW = Inches(5.8)
for i, (plan, items, color, text_color) in enumerate(plans):
    x = Inches(0.55) + i * Inches(6.55)
    y = Inches(1.55)
    is_prem = (plan == "Premium")
    add_rect(slide10, x, y, PW, Inches(5.55),
             fill_rgb=WHITE, border_rgb=color, border_pt=2.0 if is_prem else 0.75)
    add_rect(slide10, x, y, PW, Inches(0.75), fill_rgb=color)
    if is_prem:
        add_text(slide10, "⭐ RECOMANDAT",
                 x + PW - Inches(1.9), y + Inches(0.08), Inches(1.8), Inches(0.28),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide10, plan, x + Inches(0.25), y + Inches(0.18), PW - Inches(0.5), Inches(0.42),
             font_size=22, bold=True, color=WHITE)
    for j, item in enumerate(items):
        is_check = item.startswith("✓")
        clr = GREEN if is_check else RGBColor(0xEF,0x44,0x44)
        add_text(slide10, item,
                 x + Inches(0.25), y + Inches(0.95) + j * Inches(0.55),
                 PW - Inches(0.4), Inches(0.45),
                 font_size=13, color=clr)

# =============================================================================
# SLIDE 11 — CONCLUZIE / CONTACT
# =============================================================================
slide11 = prs.slides.add_slide(BLANK)
add_rect(slide11, 0, 0, W, H, fill_rgb=BLUE_DARK)

# Decorativ
for cx, cy, cr, c in [
    (Inches(12.0), Inches(0.8), Inches(3.0), RGBColor(0x1D,0x4E,0xD8)),
    (Inches(1.5),  Inches(6.5), Inches(2.0), RGBColor(0x1D,0x4E,0xD8)),
]:
    sh = slide11.shapes.add_shape(9, cx-cr/2, cy-cr/2, cr, cr)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

add_text(slide11, "EtoX Platform",
         Inches(1.0), Inches(1.0), Inches(10), Inches(1.2),
         font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide11, "Pregătirea pentru BAC devine mai simplă, mai organizată și mai motivantă.",
         Inches(1.0), Inches(2.2), Inches(11.2), Inches(0.7),
         font_size=22, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

# 4 highlights finale
highlights = [
    ("🎯", "Exerciții autentice\nfiltrate pe nevoile elevului"),
    ("👨‍🏫", "Verificare umană\nfeedback personalizat"),
    ("📊", "Monitorizare completă\npentru părinți"),
    ("🏆", "Motivare prin\ngamificare & XP"),
]
for i, (icon, text) in enumerate(highlights):
    x = Inches(0.7) + i * Inches(3.1)
    add_rect(slide11, x, Inches(3.3), Inches(2.8), Inches(1.8),
             fill_rgb=RGBColor(0x1E,0x29,0x5C), border_rgb=BLUE_MID, border_pt=0.75)
    add_text(slide11, icon, x, Inches(3.4), Inches(2.8), Inches(0.55),
             font_size=26, align=PP_ALIGN.CENTER)
    add_text(slide11, text, x, Inches(3.95), Inches(2.8), Inches(1.0),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide11, "Mulțumim!",
         Inches(1.0), Inches(5.35), Inches(11.2), Inches(0.7),
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide11, "Pentru întrebări și acces: contact@etox.ro  •  etox.ro",
         Inches(1.0), Inches(6.05), Inches(11.2), Inches(0.45),
         font_size=16, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)
add_text(slide11, "© 2026 EtoX Academy",
         Inches(1.0), Inches(6.7), Inches(11.2), Inches(0.35),
         font_size=12, color=RGBColor(0x60,0xA5,0xFA), align=PP_ALIGN.CENTER)

# =============================================================================
# SAVE
# =============================================================================
out_path = os.path.join(os.path.dirname(__file__), "EtoX_Platform_Prezentare_Parinti.pptx")
prs.save(out_path)
print(f"Salvat: {out_path}")
